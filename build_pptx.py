"""Generate a PowerPoint presentation describing the AI Music Video Studio solution."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------- palette ----------
BG      = RGBColor(0x0F, 0x12, 0x1A)   # near-black navy
PANEL   = RGBColor(0x1A, 0x20, 0x2C)
ACCENT  = RGBColor(0x6E, 0x56, 0xE8)   # indigo
ACCENT2 = RGBColor(0x22, 0xD3, 0xEE)   # cyan
ACCENT3 = RGBColor(0xF4, 0x7B, 0xB2)   # pink
TEXT    = RGBColor(0xE6, 0xE9, 0xF0)
MUTED   = RGBColor(0x9C, 0xA3, 0xAF)
GREEN   = RGBColor(0x34, 0xD3, 0x99)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def slide():
    s = prs.slides.add_slide(BLANK)
    # background
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.line.fill.background()
    bg.fill.solid(); bg.fill.fore_color.rgb = BG
    bg.shadow.inherit = False
    return s


def panel(s, left, top, width, height, fill=PANEL, line=None):
    p = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    p.adjustments[0] = 0.04
    p.fill.solid(); p.fill.fore_color.rgb = fill
    if line is None:
        p.line.color.rgb = ACCENT
        p.line.width = Pt(0.75)
    elif line == "none":
        p.line.fill.background()
    else:
        p.line.color.rgb = line; p.line.width = Pt(0.75)
    p.shadow.inherit = False
    return p


def text(s, left, top, width, height, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space_after=Pt(6), line_spacing=1.1):
    tb = s.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, r in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = space_after; p.line_spacing = line_spacing
        for content, size, bold, color in r:
            run = p.add_run(); run.text = content
            run.font.size = Pt(size); run.font.bold = bold
            run.font.color.rgb = color
            run.font.name = "Inter"
    return tb


def title_bar(s, kicker, title, accent=ACCENT):
    # accent stripe
    st = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(0.55), Inches(0.12), Inches(0.55))
    st.line.fill.background(); st.fill.solid(); st.fill.fore_color.rgb = accent; st.shadow.inherit = False
    text(s, Inches(0.85), Inches(0.5), Inches(11), Inches(0.4),
         [[(kicker.upper(), 12, True, accent)]], space_after=Pt(0))
    text(s, Inches(0.85), Inches(0.78), Inches(12), Inches(0.6),
         [[(title, 28, True, TEXT)]], space_after=Pt(0))


def bullets(s, left, top, width, height, items, size=16, color=TEXT, bullet=ACCENT2):
    tb = s.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(8); p.line_spacing = 1.15
        r = p.add_run(); r.text = "▸  "
        r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = bullet; r.font.name = "Inter"
        if isinstance(it, tuple):
            head, rest = it
            r2 = p.add_run(); r2.text = head + "  "
            r2.font.size = Pt(size); r2.font.bold = True; r2.font.color.rgb = TEXT; r2.font.name = "Inter"
            r3 = p.add_run(); r3.text = rest
            r3.font.size = Pt(size); r3.font.bold = False; r3.font.color.rgb = MUTED; r3.font.name = "Inter"
        else:
            r2 = p.add_run(); r2.text = it
            r2.font.size = Pt(size); r2.font.color.rgb = color; r2.font.name = "Inter"
    return tb


def footer(s, n):
    text(s, Inches(0.6), Inches(7.0), Inches(8), Inches(0.3),
         [[("AI Music Video Studio", 10, False, MUTED)]], space_after=Pt(0))
    text(s, Inches(11.8), Inches(7.0), Inches(1.2), Inches(0.3),
         [[(f"{n}", 10, False, MUTED)]], align=PP_ALIGN.RIGHT, space_after=Pt(0))


# ============================================================ SLIDE 1 — Cover
s = slide()
# decorative shapes
d = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.5), Inches(-2), Inches(7), Inches(7))
d.fill.solid(); d.fill.fore_color.rgb = ACCENT; d.line.fill.background()
d.fill.transparency = 0
# set transparency via XML
from pptx.oxml.ns import qn
sp = d.fill.fore_color._xFill.find(qn('a:srgbClr'))
alpha = sp.makeelement(qn('a:alpha'), {'val': '15000'}); sp.append(alpha)
d2 = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.8), Inches(4.2), Inches(4), Inches(4))
d2.fill.solid(); d2.fill.fore_color.rgb = ACCENT2; d2.line.fill.background()
sp = d2.fill.fore_color._xFill.find(qn('a:srgbClr'))
alpha = sp.makeelement(qn('a:alpha'), {'val': '12000'}); sp.append(alpha)

text(s, Inches(0.8), Inches(1.8), Inches(11), Inches(0.5),
     [[("TECHNICAL OVERVIEW", 14, True, ACCENT2)]], space_after=Pt(0))
text(s, Inches(0.8), Inches(2.3), Inches(11.5), Inches(2),
     [[("AI Music Video Studio", 54, True, TEXT)]], space_after=Pt(4))
text(s, Inches(0.8), Inches(3.5), Inches(11), Inches(0.8),
     [[("Song → beats & lyrics → scene plan → characters → per-scene", 22, False, MUTED)],
      [("image & video → ffmpeg assembly — fully automated.", 22, False, MUTED)]],
     space_after=Pt(0), line_spacing=1.2)
panel(s, Inches(0.8), Inches(5.2), Inches(11.6), Inches(1.3), fill=PANEL, line="none")
text(s, Inches(1.0), Inches(5.35), Inches(11), Inches(1.1),
     [[("FastAPI · SQLModel · SQLite   |   Next.js 14 · React Query · TypeScript · Tailwind", 13, True, TEXT)],
      [("All video / image / LLM via OpenRouter · fal.ai (optional) for word-level lyric timestamps", 13, False, MUTED)]],
     space_after=Pt(6))
footer(s, 1)

# ============================================================ SLIDE 2 — The Problem / What it does
s = slide()
title_bar(s, "Why this exists", "From a single song to a finished music video")
panel(s, Inches(0.6), Inches(1.55), Inches(6.0), Inches(5.2), fill=PANEL, line="none")
text(s, Inches(0.85), Inches(1.75), Inches(5.5), Inches(0.4),
     [[("THE PROBLEM", 13, True, ACCENT3)]], space_after=Pt(0))
bullets(s, Inches(0.85), Inches(2.25), Inches(5.5), Inches(4.3), [
    ("Manual music-video production is slow and expensive.", ""),
    "Each shot needs composition, characters, motion, and a cut to the beat.",
    "Existing AI tools give you clips, not a coherent song-length film.",
    "Continuity between shots is hard — models drift between generations.",
], size=15)
panel(s, Inches(6.85), Inches(1.55), Inches(6.0), Inches(5.2), fill=PANEL, line="none")
text(s, Inches(7.1), Inches(1.75), Inches(5.5), Inches(0.4),
     [[("THE SOLUTION", 13, True, GREEN)]], space_after=Pt(0))
bullets(s, Inches(7.1), Inches(2.25), Inches(5.5), Inches(4.3), [
    ("Upload a song → automatic audio analysis.", "BPM, key, beats, sections, lyrics."),
    ("LLM plans scenes", "aligned to beat & section windows, with prompts per scene."),
    ("Character cast", "with portrait refs reused at gen time."),
    ("Per-scene image + video", "two routes: OpenRouter I2V (default) / Seedance R2V audio-sync."),
    ("Frame-accurate chaining", "scene N+1 starts on scene N's real last frame."),
    ("ffmpeg assembly", "concat + song mux → streaming-ready MP4."),
], size=14, bullet=GREEN)
footer(s, 2)

# ============================================================ SLIDE 3 — Architecture / Stack
s = slide()
title_bar(s, "Architecture", "Two-service stack, one SQLite file")
# backend panel
panel(s, Inches(0.6), Inches(1.6), Inches(6.0), Inches(5.1), fill=PANEL, line="none")
text(s, Inches(0.85), Inches(1.8), Inches(5.5), Inches(0.4), [[("BACKEND  ·  FastAPI + SQLModel + SQLite", 13, True, ACCENT2)]], space_after=Pt(0))
bullets(s, Inches(0.85), Inches(2.35), Inches(5.5), Inches(4.2), [
    ("routers/", "projects · songs · scenes · generation"),
    ("services/openrouter.py", "video, image, LLM, transcription"),
    ("scene_planner.py", "batch plan · expand · continuation · characters"),
    ("generation_service.py", "per-scene image→video pipeline"),
    ("assembly.py", "ffmpeg concat + audio mux + faststart"),
    ("versioning.py / urls.py / llm_json.py", "shared invariants + helpers"),
], size=13)
# frontend panel
panel(s, Inches(6.85), Inches(1.6), Inches(6.0), Inches(5.1), fill=PANEL, line="none")
text(s, Inches(7.1), Inches(1.8), Inches(5.5), Inches(0.4), [[("FRONTEND  ·  Next.js 14 + React Query", 13, True, ACCENT3)]], space_after=Pt(0))
bullets(s, Inches(7.1), Inches(2.35), Inches(5.5), Inches(4.2), [
    ("FlowStudio", "per-project workspace orchestrating cells"),
    ("cells/", "Song · Characters · Plan · Generate · Assemble"),
    ("SceneGenRow", "per-scene row: chain 🔗, wand ✨, mic toggle"),
    ("lib/api.ts", "typed fetch wrapper, extracts FastAPI detail"),
    ("useConfirm()", "in-app overlay replaces window.confirm()"),
    ("Auto-polling", "plan 2s · assemble 3s · portrait 2.5s"),
], size=13, bullet=ACCENT3)
footer(s, 3)

# ============================================================ SLIDE 4 — Pipeline flow
s = slide()
title_bar(s, "Pipeline", "End-to-end flow, v1")
flow = [
    ("1", "Upload song",          "librosa → BPM, beats, sections · whisper → lyrics · LLM → theme", ACCENT2),
    ("2", "Suggest characters",   "scene_planner proposes a cast from the song's theme", ACCENT2),
    ("3", "Plan scenes (batch)",  "LLM divides song into shots aligned to beat windows", ACCENT),
    ("4", "Render image per scene", "OpenRouter image gen + character refs", ACCENT3),
    ("5", "Render video per scene", "I2V default  /  Seedance R2V audio-sync (opt-in)", GREEN),
    ("6", "Chain + wand",          "vision-LLM writes next prompt from real last frame", GREEN),
    ("7", "Assemble final MP4",    "ffmpeg concat + song mux + -movflags +faststart", ACCENT3),
]
top = Inches(1.55); h = Inches(0.78); gap = Inches(0.12)
for i, (n, head, sub, col) in enumerate(flow):
    t = Emu(top + i * (h + gap))
    panel(s, Inches(0.6), t, Inches(12.1), h, fill=PANEL, line="none")
    # number chip
    chip = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.78), t + Inches(0.16), Inches(0.46), Inches(0.46))
    chip.fill.solid(); chip.fill.fore_color.rgb = col; chip.line.fill.background(); chip.shadow.inherit = False
    tf = chip.text_frame; tf.text = n
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].runs[0].font.bold = True; tf.paragraphs[0].runs[0].font.size = Pt(15); tf.paragraphs[0].runs[0].font.color.rgb = BG
    text(s, Inches(1.45), t + Inches(0.12), Inches(4.0), Inches(0.55),
         [[(head, 16, True, TEXT)]], anchor=MSO_ANCHOR.MIDDLE, space_after=Pt(0))
    text(s, Inches(5.6), t + Inches(0.12), Inches(7.0), Inches(0.55),
         [[(sub, 13, False, MUTED)]], anchor=MSO_ANCHOR.MIDDLE, space_after=Pt(0))
footer(s, 4)

# ============================================================ SLIDE 5 — Two video routes
s = slide()
title_bar(s, "Two video routes", "Picked per scene, by model + mic toggle")
# left route
panel(s, Inches(0.6), Inches(1.6), Inches(6.0), Inches(5.1), fill=PANEL, line="none")
text(s, Inches(0.85), Inches(1.8), Inches(5.5), Inches(0.4), [[("ROUTE A — OpenRouter Image-to-Video  (default)", 13, True, ACCENT2)]], space_after=Pt(0))
text(s, Inches(0.85), Inches(2.2), Inches(5.5), Inches(0.4), [[("Cheap · no audio · ships first_frame + optional refs", 13, False, MUTED)]], space_after=Pt(0))
bullets(s, Inches(0.85), Inches(2.8), Inches(5.5), Inches(3.8), [
    ("first_frame =", "chained prev's last frame, else scene's still"),
    ("input_references =", "character portraits (Seedance only; Kling/Veo drop)"),
    ("Supported models:", "Seedance, Kling, Veo"),
    ("Tradeoff:", "lowest cost; no character lip-sync"),
], size=13)
# right route
panel(s, Inches(6.85), Inches(1.6), Inches(6.0), Inches(5.1), fill=PANEL, line="none")
text(s, Inches(7.1), Inches(1.8), Inches(5.5), Inches(0.4), [[("ROUTE B — fal Seedance Reference-to-Video  (opt-in)", 13, True, GREEN)]], space_after=Pt(0))
text(s, Inches(7.1), Inches(2.2), Inches(5.5), Inches(0.4), [[("~6× cost · audio in · character “performs” with lipsync", 13, False, MUTED)]], space_after=Pt(0))
bullets(s, Inches(7.1), Inches(2.8), Inches(5.5), Inches(3.8), [
    ("image_urls =", "first-frame source + up to 9 char refs"),
    ("No first_frame concept", "every image is a reference"),
    ("audio_sync_enabled +", "model.supports_audio_input"),
    ("_extract_audio_segment", "slices song[start..end] as synthesis input"),
], size=13, bullet=GREEN)
footer(s, 5)

# ============================================================ SLIDE 6 — Scene chaining
s = slide()
title_bar(s, "Continuity", "Frame-accurate chaining between scenes")
# old vs new
panel(s, Inches(0.6), Inches(1.6), Inches(6.0), Inches(2.35), fill=PANEL, line=ACCENT3)
text(s, Inches(0.85), Inches(1.75), Inches(5.5), Inches(0.4), [[("OLD — removed", 13, True, ACCENT3)]], space_after=Pt(0))
bullets(s, Inches(0.85), Inches(2.2), Inches(5.5), Inches(1.7), [
    "Passed scene N+1's planned still as a soft target.",
    "Model rarely landed on it → visible seam.",
], size=13, bullet=ACCENT3)
panel(s, Inches(6.85), Inches(1.6), Inches(6.0), Inches(2.35), fill=PANEL, line=GREEN)
text(s, Inches(7.1), Inches(1.75), Inches(5.5), Inches(0.4), [[("NEW — current", 13, True, GREEN)]], space_after=Pt(0))
bullets(s, Inches(7.1), Inches(2.2), Inches(5.5), Inches(1.7), [
    "chain_from_prev=true → scene N+1 first_frame = N's extracted last frame.",
    "Real rendered pixels → pixel-accurate handoff at the seam.",
], size=13, bullet=GREEN)
# iterative workflow
panel(s, Inches(0.6), Inches(4.2), Inches(12.1), Inches(2.5), fill=PANEL, line="none")
text(s, Inches(0.85), Inches(4.35), Inches(11.5), Inches(0.4), [[("ITERATIVE BUILD WORKFLOW", 13, True, ACCENT2)]], space_after=Pt(0))
bullets(s, Inches(0.85), Inches(4.85), Inches(11.5), Inches(1.8), [
    ("Plan", "“Just scene 1” → one scene from the seed."),
    ("Render", "Img then Vid on scene 1 → last frame extracted automatically."),
    ("Chain", "click 🔗 → empty scene 2 chained from scene 1."),
    ("Wand", "click ✨ → vision-LLM reads scene 1's last frame, writes scene 2's prompt."),
    ("Repeat", "Vid → 🔗 → ✨ → Vid, until the song ends."),
], size=13)
footer(s, 6)

# ============================================================ SLIDE 7 — Engineering conventions
s = slide()
title_bar(s, "Engineering", "Conventions that keep the codebase honest")
two_col = [
    ("Storage URLs", "to_storage_url() — never hardcode localhost:8010", ACCENT2),
    ("LLM JSON", "parse_llm_json() tolerates fences, prose, singleton arrays", ACCENT2),
    ("Versioned rows", "make_active() / delete_and_promote() enforce one-active invariant", ACCENT2),
    ("Prompt history", "_save_prompt_version() — Scene.image_prompt is a mirror only", ACCENT2),
    ("Endpoints", "top-level try/except + traceback → specific 500 message", ACCENT3),
    ("Migrations", "_apply_schema_migrations() ALTERs tables for new columns", ACCENT3),
    ("Confirms", "useConfirm() overlay — never window.confirm()", ACCENT3),
    ("Errors → UI", "request() extracts FastAPI detail; no re-parsing", ACCENT3),
    ("Avoid", "no post-process lipsync, no last_frame chaining, no motion-offset trim", ACCENT3),
]
# render as 3x3 grid of cards
cw, ch = Inches(4.0), Inches(1.55)
x0, y0 = Inches(0.6), Inches(1.6)
gx, gy = Inches(0.15), Inches(0.15)
for i, (head, body, col) in enumerate(two_col):
    r, c = divmod(i, 3)
    x = x0 + c * (cw + gx); y = y0 + r * (ch + gy)
    panel(s, x, y, cw, ch, fill=PANEL, line="none")
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.18), y + Inches(0.2), Inches(0.18), Inches(0.18))
    dot.fill.solid(); dot.fill.fore_color.rgb = col; dot.line.fill.background(); dot.shadow.inherit = False
    text(s, x + Inches(0.45), y + Inches(0.12), cw - Inches(0.5), Inches(0.4),
         [[(head, 14, True, TEXT)]], space_after=Pt(0))
    text(s, x + Inches(0.2), y + Inches(0.55), cw - Inches(0.35), Inches(0.95),
         [[(body, 11, False, MUTED)]], space_after=Pt(0), line_spacing=1.1)
footer(s, 7)

# ============================================================ SLIDE 8 — Known gotchas / tradeoffs
s = slide()
title_bar(s, "Gotchas & tradeoffs", "Provider drift and model constraints")
items = [
    ("Seedance content filter", "refuses photoreal portraits as input_references — surfaced explicitly, no silent degrade.", ACCENT3),
    ("Veo needs allow_adult", "personGeneration='allow_adult' set on every Veo submission or it refuses.", ACCENT3),
    ("Veo + refs incompatible", "ignores input_references when first_frame present → supports_reference_images=False.", ACCENT3),
    ("Kling drops refs on OR", "input_references exposed natively but OpenRouter passthrough drops them.", ACCENT3),
    ("Gemini Image", "sometimes returns text, not an image — auto-retry-on-text handles it.", ACCENT2),
    ("Backend reload kills in-flight", "uvicorn --reload; mitigated by graceful shutdown + short LLM batches + frontend retry.", ACCENT2),
    ("Chain needs prev rendered", "chain_from_prev reads extracted_last_frame_path; generate sequentially.", ACCENT2),
    ("Durations snapped to model", "10s scene on Veo Lite → 8s; use Kling/Seedance for >8s scenes.", ACCENT2),
    ("last frame mutated in place", "to_storage_url(path, cache_bust=True) breaks stale browser cache.", ACCENT2),
]
top = Inches(1.55); h = Inches(0.56); gap = Inches(0.07)
for i, (head, body, col) in enumerate(items):
    t = Emu(top + i * (h + gap))
    panel(s, Inches(0.6), t, Inches(12.1), h, fill=PANEL, line="none")
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), t, Inches(0.08), h)
    bar.fill.solid(); bar.fill.fore_color.rgb = col; bar.line.fill.background(); bar.shadow.inherit = False
    text(s, Inches(0.85), t, Inches(3.6), h,
         [[(head, 13, True, TEXT)]], anchor=MSO_ANCHOR.MIDDLE, space_after=Pt(0))
    text(s, Inches(4.5), t, Inches(8.0), h,
         [[(body, 12, False, MUTED)]], anchor=MSO_ANCHOR.MIDDLE, space_after=Pt(0))
footer(s, 8)

# ============================================================ SLIDE 9 — Config / models
s = slide()
title_bar(s, "Configuration", "config.py — models, pricing, capabilities")
panel(s, Inches(0.6), Inches(1.6), Inches(4.0), Inches(5.1), fill=PANEL, line="none")
text(s, Inches(0.85), Inches(1.8), Inches(3.5), Inches(0.4), [[("VIDEO_MODELS", 13, True, ACCENT2)]], space_after=Pt(0))
bullets(s, Inches(0.85), Inches(2.35), Inches(3.5), Inches(4.2), [
    "Seedance 1.0 / 2.0 (refs ✓, audio ✓)",
    "Kling 1.6 (refs ✗ on OR)",
    "Veo 3 / Veo 3 Fast (refs ✗)",
    "Each: durations, resolutions, pricing",
], size=12)
panel(s, Inches(4.7), Inches(1.6), Inches(4.0), Inches(5.1), fill=PANEL, line="none")
text(s, Inches(4.95), Inches(1.8), Inches(3.5), Inches(0.4), [[("IMAGE_MODELS", 13, True, ACCENT3)]], space_after=Pt(0))
bullets(s, Inches(4.95), Inches(2.35), Inches(3.5), Inches(4.2), [
    "Gemini Image (default)",
    "Recraft / others selectable",
    "Retry-on-text built in",
    "Celebrity names trip likeness filter",
], size=12, bullet=ACCENT3)
panel(s, Inches(8.8), Inches(1.6), Inches(4.0), Inches(5.1), fill=PANEL, line="none")
text(s, Inches(9.05), Inches(1.8), Inches(3.5), Inches(0.4), [[("ENV VARS", 13, True, GREEN)]], space_after=Pt(0))
bullets(s, Inches(9.05), Inches(2.35), Inches(3.5), Inches(4.2), [
    ("Required:", "OPENROUTER_API_KEY"),
    ("Optional:", "FAL_API_KEY (word-level lyrics)"),
    ("Optional:", "SUNO_API_KEY (music gen)"),
    ("Override:", "PUBLIC_BASE_URL, STORAGE_DIR, port 8010"),
], size=12, bullet=GREEN)
footer(s, 9)

# ============================================================ SLIDE 10 — Closing
s = slide()
d = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(3), Inches(6), Inches(6))
d.fill.solid(); d.fill.fore_color.rgb = ACCENT2; d.line.fill.background()
sp = d.fill.fore_color._xFill.find(qn('a:srgbClr')); alpha = sp.makeelement(qn('a:alpha'), {'val':'10000'}); sp.append(alpha)
text(s, Inches(0.8), Inches(2.2), Inches(11.5), Inches(0.5), [[("SUMMARY", 14, True, ACCENT2)]], space_after=Pt(0))
text(s, Inches(0.8), Inches(2.7), Inches(11.5), Inches(1.5),
     [[("An end-to-end AI music-video studio.", 38, True, TEXT)]], space_after=Pt(6), line_spacing=1.15)
bullets(s, Inches(0.8), Inches(4.2), Inches(11.5), Inches(2.5), [
    ("One SQLite file, two services", "FastAPI backend + Next.js frontend, all AI via OpenRouter."),
    ("Two video routes per scene", "cheap I2V default, Seedance R2V audio-sync opt-in."),
    ("Pixel-accurate continuity", "chained scenes start on the previous clip's real last frame."),
    ("Convention-driven codebase", "shared helpers enforce invariants; gotchas documented up front."),
], size=16)
text(s, Inches(0.8), Inches(6.7), Inches(11.5), Inches(0.4),
     [[("Backend :8010   ·   Frontend :3000   ·   API docs :8010/docs", 12, False, MUTED)]], space_after=Pt(0))
footer(s, 10)

out = "AI-Music-Video-Studio.pptx"
prs.save(out)
print("wrote", out)